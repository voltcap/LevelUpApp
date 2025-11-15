import os
import json
import re
from pathlib import Path
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import docx
import PyPDF2
import openai

openai.api_key = #add your api key
openai.api_base = "https://openrouter.ai/api/v1"
MODEL_NAME = "tngtech/deepseek-r1t2-chimera:free"

app = FastAPI()

from fastapi.middleware.cors import CORSMiddleware

origins = [
    "http://localhost:3000",
    "http://127.0.0.1:3000",
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins, 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def readPdf(f):
    text = ""
    reader = PyPDF2.PdfReader(f)
    for page in reader.pages:
        content = page.getText()
        if content:
            text += content + "\n"
    return text

def readWord(f):
    doc = docx.Document(f)
    return "\n".join([p.text for p in doc.paragraphs])

def readPpt(f):
    prs = Presentation(f)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def readTxt(f):
    return f.read().decode("utf-8")

def loadDocs(file: UploadFile):
    ext = Path(file.filename).suffix.lower()
    if ext == ".pdf":
        return readPdf(file.file)
    elif ext == ".docx":
        return readWord(file.file)
    elif ext == ".pptx":
        return readPpt(file.file)
    elif ext in [".txt", ".md"]:
        return readTxt(file.file)
    elif ext in [".py", ".js", ".java", ".cpp", ".c", ".html", ".css"]:
        return readTxt(file.file)
    else:
        raise ValueError(f"upload supported files {ext}")

def chimeraModel(prompt):
    try:
        response = openai.ChatCompletion.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": "You generate quizzes in strict JSON."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        return response["choices"][0]["message"]["content"]
    except Exception as e:
        return f"API Error: {e}"

def levelUpQuiz(full_text, difficulty, quiz_type, num_questions):
    prompt = f"""
Generate a quiz from the following content.

CONTENT:
\"\"\"{full_text[:12000]}\"\"\"

REQUIREMENTS:
- Difficulty: {difficulty}
- Quiz format: {quiz_type} (multiple-choice, true-false, fill-in-the-blank)
- Number of questions: EXACTLY {num_questions}

Return ONLY valid JSON in this exact format:
{{
  "quiz": [
    {{
      "question": "Your question here",
      "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
      "answer_index": 0
    }}
  ]
}}
"""
    raw = chimeraModel(prompt)

    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not match:
        return {"error": "Model did not return JSON", "raw": raw}

    try:
        return json.loads(match.group(0))
    except Exception as e:
        return {"error": f"JSON parsing failed: {e}", "raw": raw}

@app.post("/levelUpQuiz")
async def api_levelUpQuiz(
    file: UploadFile = File(...),
    difficulty: str = "easy",
    quiz_type: str = "mcq",
    num_questions: int = 5
):
    try:
        full_text = loadDocs(file)
    except Exception as e:
        return {"error": f"Failed to read document: {e}"}

    quiz = levelUpQuiz(full_text, difficulty, quiz_type, num_questions)
    return quiz
